import json
import os
from pptx import Presentation
import re
import json
import pandas as pd
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime

## Set correct working directory
script_dir = os.path.dirname(os.path.abspath(__file__))
os.chdir(script_dir)


########################################################################################################################
###
### Create variables
###
########################################################################################################################

## Creates the American to British ESV Dictionary
list1 = [
         "(\W)([Ss])avior(s)?(\W)",
         "(\W)neighbor(s|ing)?(\W)",
         "(\W)favor(able|ite|s|ed|itism)?(\W)",
         "(\W)Favor(\W)",
         "(\W)labor(ed|s)?(\W)",
         "(\W)(vap|vig)or(\W)",
         "(\W)clamor(\W)",
         "(\W)([Ss])plendor(\W)",
         "(\W)color(s|ed)?(\W)",
         "(\W)([Hh])onor(s|able|ing|ed)?(\W)",
         "(\W)dishonor(s|able|ing|ed)?(\W)",
         "(\W)travel(ed|er|ers|ing)(\W)",
         "(\W)marvel(ous|ously|ed|ing)(\W)",
         "(\W)([Cc])ounsel(or|ors|ed)(\W)",
         "(\W)plow(s|ed|ers|ing|man|men|share|shares)?(\W)",
         "(\W)judgment(s)?(\W)",
         "(\W)(recogn|[Rr]eal|[Oo]rgan|[Ss]ymbol|bapt|critic|apolog|sympath)iz(e|ed|es|ing)?(\W)",
         "(\W)(un)?authorized(\W)",
         "(\W)(centi)?meters(\W)",
         "(\W)liter(s)?(\W)",
         "(\W)scepter(s)?(\W)",
         "(\W)worship(ed|er|ers|ing)(\W)",
         "(\W)quarrel(ed|ing)(\W)",
         "(\W)benefited(\W)",
         "(\W)signaled(\W)",
         "(\W)paralyzed(\W)",
         "(\W)fulfill(s|ment)?(\W)",
         "(\W)skillful(ly)?(\W)",
         "(\W)jewelry(\W)",
         "(\W)(De|de|of)fense(s|less)?(\W)",
         "(\W)([Ss])ulfur?(\W)"
    ]

list2 = [
        r'\1\2aviour\3\4',
        r'\1neighbour\2\3',
        r'\1favour\2\3',
        r'\1Favour\2',
        r'\1labour\2\3',
        r'\1\2our\3',
        r'\1clamour\2',
        r'\1\2plendour\3',
        r'\1colour\2\3',
        r'\1\2onour\3\4',
        r'\1dishonour\2\3',
        r'\1travell\2\3',
        r'\1marvell\2\3',
        r'\1\2ounsell\3\4',
        r'\1plough\2\3',
        r'\1judgement\2\3',
        r'\1\2is\3\4',
        r'\1\2authorised\3',
        r'\1\2metres\3',
        r'\1litre\2\3',
        r'\1sceptre\2\3',
        r'\1worshipp\2\3',
        r'\1quarrell\2\3',
        r'\1benefitted\2',
        r'\1signalled\2',
        r'\1paralysed\2',
        r'\1fulfil\2\3',
        r'\1skilful\2\3',
        r'\1jewellery\2',
        r'\1\2fence\3\4',
        r'\1\2ulphur\3'
    ]
#############################################################

## Generates list to use in below function
anglo_list = list(zip(list1, list2))

# If value is a list, then use second value as slide title
slide_dict = {
    2: "notices",
    3: ["call_to_worship","Call to worship"],
    4: "song1",
    5: ["confession","Confession"],
    6: ["assurance", "Assurance of Pardon"],
    7: ["lords_prayer","The Lord's Prayer"],
    8: "psalm",
    9: ["first_reading","First Reading"],
    10: "catechism_reading_previous",
    11: "catechism_reading_today",
    12: "song2",
    13: ["second_reading","Second Reading"],
    14: ["apostles_creed","Apostles' Creed (180 AD)"],
    15: ["prayers","Prayers of Intercession"],
    16: "song3",
    17: "song4",
    18: ["the_grace","The Grace"],
    19: "goodbye"
}

# Make a dictionary of the keys where slide_dict value starts with "song" or "psalm", some values are lists, so need to check for that
song_list = []
component_list = []
reading_list = []
catechism_list = []
for elem in slide_dict:
    if type(slide_dict[elem]) != list:
        if slide_dict[elem].startswith("song") or slide_dict[elem].startswith("psalm"):
            song_list.append(elem)
        elif slide_dict[elem].startswith("catechism"):
            catechism_list.append(elem)
    elif type(slide_dict[elem]) == list:
        if "reading" in slide_dict[elem][0]:
            reading_list.append(elem)
        elif slide_dict[elem][0] not in ['call_to_worship'] and slide_dict[elem][0] not in ['prayers']:
            component_list.append(elem)

### Load in json files
with open("psalms.json", "r") as f:
    psalms_json = json.load(f)

with open("wsc.json", "r") as f:
    wsc_json = json.load(f)

with open("components.json", "r") as f:
    components_json = json.load(f)

## ESV API keys
API_KEY = 'ade14fe748fbb522b8dfb225ec6b222fa148cddc'
API_URL = 'https://api.esv.org/v3/passage/text/'

########################################################################################################################
####
#### Functions
####
########################################################################################################################

## Used to get the readings from the google sheet
online_csv = 'https://docs.google.com/spreadsheets/d/e/2PACX-1vSIEtCzAWNJVZK7T1OEo1oGhnTib2bgFfdYRfFON1gpbG7LkHTFRJKSioV087Ys1oBZci80cRIRm0_u/pub?gid=0&single=true&output=csv'

## Extracts the relevant details from the components JSON
def component_assigner(val, flag = None):
    if flag != 'communion':
        res = list(filter(lambda item: item["Component"] == val, components_json))
    else:
        initial = list(filter(lambda item: item["Component"] == "communion", components_json))
        res = list(filter(lambda item: item["Component"] == val, initial[0]["Content"]))
        
    speaker = res[0]["Content"]["speaker"]
    return_list = []
    for i in range(1,5):
        try:
            return_list.append(res[0]["Content"][str(i)])
        except:
            pass
    return speaker, return_list
'''
Returns the speaker, which most likely goes on every slide
And returns the text for each slide, in a list, with the content for each slide in each element
'''

## Function to extract the appropriate catechism question
def catechism_finder(val):
    # if val is not int cast to int
    if type(val) != int:
        val = int(val)
    real_val = val - 1
    question = wsc_json['Data'][real_val]["Question"]
    answer = wsc_json['Data'][real_val]["Answer"]
    return (question,answer, val)


## Function to convert American language text to British
def replace_text(text, rep_text=anglo_list):
    for pattern, replacement in rep_text:
        pattern = re.compile(pattern)
        text = pattern.sub(replacement, text)
    return text

## Function to convert digits in verse to superscript
def super_verse(text):
    pattern = re.compile(r'\[([^\[\]]*)\]\s')
    matches = pattern.findall(text)
    for match in matches:
        superscript = ''.join(['⁰¹²³⁴⁵⁶⁷⁸⁹'[int(x)] if x.isdigit() else x for x in match])
        text = text.replace('[' + match + '] ', superscript)
    return text
    

## ESV API to return text from passage text
def get_esv_text(passage):
  params = {
    'q': passage,
    'include-headings': False,
    'include-footnotes': False,
    'include-verse-numbers': True,
    'include-short-copyright': False,
    'include-passage-references': False
  }
  
  headers = {
    'Authorization': 'Token %s' % API_KEY
  }
  
  response = requests.get(API_URL, params=params, headers=headers)
  
  passages = response.json()['passages']

  result = passages[0].strip()
  result = super_verse(result)
  result = replace_text(result)
  
  if passages:
    return (result,passage)
  else:
    return 'Error: Passage not found'

## Hymn Scraper Class
class HymnScraper:
    def __init__(self, url1 = None, url2 = None):
        self.url1 = url1
        self.url2 = url2
        self.title = None
        self.author = None
        self.copyr = None
        self.n_verses = None
        self.composer = None
        self.tune = None
        self.meter = None
        
    def hymn_lyrics(self):
        ## Word Extraction
        response = requests.get(self.url1)
        soup = BeautifulSoup(response.content, 'html.parser')
        try:
            self.title = soup.find("span", class_="hy_infoLabel", string="Title:").find_next().text
        except:
            None
        try:
            self.author = soup.find("span", class_="hy_infoLabel", string="Author:").find_next().text
        except:
            try:
                self.author = soup.find("span", class_="hy_infoLabel", string="Author (attributed to):").find_next().text
                self.author = self.author + " (atrb)"
            except:
                self.author = None
        try:
            self.copyr = soup.find("span", class_="hy_infoLabel", string="Copyright:").find_next().text
        except:
            self.copyr = None

        try:
            body_text = soup.select_one('div#at_fulltext.authority_section div div.authority_columns')
            paragraphs = body_text.find_all('p')
            verses = []
            for p in paragraphs:
                verses.append(p.text)

            self.n_verses = []
            for item in verses:
                if item[0].isdigit():
                    items = item.split(" ")
                    item = " ".join(items[1:])
                    if item.rstrip().endswith(" [Refrain]"):
                        item = item.split(" [Refrain]")[0]
                        self.n_verses.append(item)
                        self.n_verses.append(refrain)
                    else:
                        self.n_verses.append(item)
                if item.split(":")[0] == 'Refrain':
                    refrain = item.split(":")[1] + "\n"
                    self.n_verses.append(refrain)
        except:
            self.n_verses = None

    def tune_details(self):
        ## Tune Extraction
        response = requests.get(self.url2)
        soup = BeautifulSoup(response.content, 'html.parser')
        try:
            self.composer = soup.find("span", class_="hy_infoLabel", string="Composer:").find_next().text
        except:
            try:
                self.composer = soup.find("span", class_="hy_infoLabel", string="Composer (attributed to):").find_next().text
                self.composer = self.composer + " (atrb)"
            except:
                self.composer = "Unknown"
        try:
            self.tune = soup.find("span", class_="hy_infoLabel", string="Title:").find_next().text
        except:
            self.tune = None
        try:
            self.meter = soup.find("span", class_="hy_infoLabel", string="Meter:").find_next().text
        except:
            self.meter = None


    def get_lyrics(self):
        return (self.n_verses,self.title, self.author, self.composer, self.tune, self.copyr)
    
    def get_tune(self):
        return (self.composer, self.tune, self.meter)

## Hymn Scraper Functions
## Function to instantiate class and extract details
def song_details(url1,url2):
    scraper = HymnScraper(url1, url2)
    scraper.hymn_lyrics()
    scraper.tune_details()
    return scraper.get_lyrics()

def tune_details(url2):
    scraper = HymnScraper(None, url2)
    scraper.tune_details()
    return scraper.get_tune()


# function to convert to superscript
def get_super(x):
    normal = "0123456789-"
    super_s = "⁰¹²³⁴⁵⁶⁷⁸⁹⁻"
    res = x.maketrans(''.join(normal), ''.join(super_s))
    return x.translate(res)

def super_psalm(text):
    # Regular expression to match numbers or numbers with dashes between, but not if they have a space on both sides
    pattern = r"(\d+\-?\d*)(?=[^\s\d])"
    def get_super(x):
        normal = "0123456789-"
        super_s = "⁰¹²³⁴⁵⁶⁷⁸⁹⁻"
        res = x.maketrans(''.join(normal), ''.join(super_s))
        return x.translate(res)
    # Use re.sub to replace the matched numbers with superscript format
    text = re.sub(pattern, lambda x: get_super(x.group()), text)
    return text


def psalm_getter(psalm_name, psalms_json=psalms_json):
    
    match = re.search(r"^(Psalm )?(\d{1,3})(:(([1-9]\d{0,2})-([1-9]\d{0,2})))?(\s\(([a-zA-Z])\))?(\s\((\d{0,2})\))?$", psalm_name)

    if match:
        if match.group(2):
            psalm = match.group(2)
        if match.group(7):
            version = match.group(8)
        else:
            version = "a"
        if match.group(9):
            section = match.group(10)
        else:
            section = None
        if match.group(4):
            verses = match.group(4)
        elif section != None:
            verses = "1-300"
        else:
            verses = "1-30"

    if section == None:
        body = [d for d in psalms_json if d["Psalm"] == str(psalm) and d["Content"]["Version"] == version][0]['Content']['Body']
        meter = [d for d in psalms_json if d["Psalm"] == str(psalm) and d["Content"]["Version"] == version][0]['Content']['Meter']
    else:
        body = [d for d in psalms_json if d["Psalm"] == str(psalm) and d["Content"]["Section"] == str(section)][0]['Content']['Body']
        meter = [d for d in psalms_json if d["Psalm"] == str(psalm) and d["Content"]["Section"] == str(section)][0]['Content']['Meter']

    lst = []
    rng = [int(n) for n in verses.split("-")]
    for ints in range(rng[0],rng[1]):
        try:
            lst.append(body[str(ints)])
        except:
            pass

    lst = list(map(super_psalm,lst))

    return (psalm_name, meter,lst)

'''
# Pass in the string of the Psalm in the below format
##string = "Psalm 19:2-10 (a)"
# Swapping a either for section (number) or version (letter)
psalm, meter, lst = psalm_getter(string, psalms_json)
print(psalm,meter,lst)
'''
## Creates all the variables that we want to use in the template
## Using the dataframe we feed in the cycle number
df = pd.read_csv(online_csv)

cycle = df["Example Column"]

try:
    call_to_worship = get_esv_text(cycle[1])
    song1 = song_details(cycle[2], cycle[3])
    ## (psalm, meter,lst)
    psalm = psalm_getter(cycle[4])
    psalm_tune = tune_details(cycle[5])
    first_reading = cycle[6]
    catechism_reading_today = catechism_finder(cycle[7])
    catechism_reading_previous = catechism_finder(int(cycle[7])-1)
    song2 = song_details(cycle[8], cycle[9])
    second_reading = cycle[10]
    song3 = song_details(cycle[11], cycle[12])
    try:
        song4 = song_details(cycle[13], cycle[14])
    except:
        pass
except:
    pass

class SlideFill:
    def __init__(self, slide):
        self.slide = slide
    
    def fill_main(self, title, body, copy):
        for idx, placeholder in enumerate(self.slide.placeholders):
            try:
                placeholder.text_frame.clear()
                if idx == 0:
                    new_body = title
                elif idx == 1:
                    new_body = body
                elif idx == 2:
                    new_body = copy
                placeholder.text_frame.text = new_body
            except:
                pass
    
    def fill_component(self, title, body, address):
        for idx, placeholder in enumerate(self.slide.placeholders):
            try:
                placeholder.text_frame.clear()
                if idx == 0:
                    placeholder.text_frame.text = title
                elif idx == 1:
                    p = placeholder.text_frame.add_paragraph()
                    run = p.add_run()
                    run.font.size = Pt(28)
                    run.font.color.rgb = RGBColor(161, 38, 38)
                    run.text = f"{address} "
                    run = p.add_run()
                    run.font.size = Pt(28)
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    run.text = body    
            except:
                pass
    
    def fill_psalm(self, title, body, copy, meter):
        for idx, placeholder in enumerate(self.slide.placeholders):
            try:
                placeholder.text_frame.clear()
                if idx == 0:
                    new_body = title
                elif idx == 1:
                    new_body = body
                elif idx == 2:
                    new_body = copy
                elif idx == 3:
                    new_body = meter
                placeholder.text_frame.text = new_body
            except:
                pass

    def fill_reading(self,title,body):
        for idx, placeholder in enumerate(self.slide.placeholders):
            try:
                placeholder.text_frame.clear()
                if idx == 0:
                    new_body = title
                elif idx == 1:
                    new_body = body
                placeholder.text_frame.text = new_body
            except:
                pass

############################################################################################################

def flag_format(flag):
    try:
        ## Such as 'notices' and 'song1'
        if type(flag) == str:
            return globals()[flag]
        elif type(flag) == list:
            try:
                ## For 'component versions
                if type(flag) == list:
                    return component_assigner(flag[0])
            except:
                ## For 'call to worship'
                return (globals()[flag[0]],flag[1])
    except:
        return "Not found"

def slide_maker(layout_type, prs):
    slide_layout = prs.slide_layouts[layout_type]
    slide = prs.slides.add_slide(slide_layout)
    return slide


def slide_writer(flag, prs):
    # Populate the placeholders on the slide with data from variables
    flag_val = flag_format(slide_dict[flag])

    
    if flag == 3:
        slide = slide_maker(1, prs)
        # This works for call to worship
        title = flag_val[1]
        body = flag_val[0][0]
        copy = flag_val[0][1]
        SlideFill(slide).fill_main(title, body, copy)
    
    elif slide_dict[flag] == "goodbye":
        slide = slide_maker(7, prs)

    elif flag in song_list and flag_val != "Not found":
        # This works for psalms
        try:
            flag_val[0].find("Psalm") != -1
            list_len = len(flag_val[2])
            for i in range(0,list_len):
                slide = slide_maker(3, prs)
                title = flag_val[0]
                body = flag_val[2][i]
                copy = f"Words: Sing Psalms! © 2003 Free Church of Scotland\nComposer: {psalm_tune[0]}\nTune: {psalm_tune[1]}\n©: Public Domain\nCCLI: 522221"
                meter = f"Meter: {flag_val[1]}"
                SlideFill(slide).fill_psalm(title, body, copy, meter)
        except:
            # This works for songs
            try:
                list_len = len(flag_val[0])
                for i in range(0,list_len):
                    slide = slide_maker(0, prs)
                    title = flag_val[1]
                    body = flag_val[0][i]
                    # Could put this outside the loop, in the future to only fill on the final slide
                    copy = f"Words: {flag_val[2]}\nComposer: {flag_val[3]}\nTune: {flag_val[4]}\n©: {flag_val[5]}\nCCLI: 522221"
                    SlideFill(slide).fill_main(title, body, copy)
            except:
                slide = slide_maker(0, prs)
                try:
                    title = flag_val[1]
                except:
                    title = slide_dict[flag]
                body = "Not in Public Domain"
                try:
                    copy = f"Words: {flag_val[2]}\nComposer: {flag_val[3]}\nTune: {flag_val[4]}\n©: {flag_val[5]}\nCCLI: 522221"
                except:
                    copy = "Error"
                SlideFill(slide).fill_main(title, body, copy)
                
    elif flag in component_list:
        # This works for components
        list_len = len(flag_val[1])
        for i in range(0,list_len):
            slide = slide_maker(2, prs)
            title = slide_dict[flag][1]
            address = flag_val[0]
            body = flag_val[1][i]
            SlideFill(slide).fill_component(title, body, address)
    
    elif flag in reading_list:
        title = flag_val[1]
        body = f"{flag_val[0]}\n\npg. X"
        slide = slide_maker(4, prs)
        SlideFill(slide).fill_reading(title, body)

    elif flag in catechism_list:
        title = f"Westminster Shorter Catechism {flag_val[2]}"
        question = f"\n\n{flag_val[1]}"
        answer = f"{flag_val[0]}"
        slide = slide_maker(5, prs)
        SlideFill(slide).fill_component(title, question, answer)

    elif "prayer" in slide_dict[flag][0]:
        title = slide_dict[flag][1]
        body = ""
        slide = slide_maker(4, prs)
        SlideFill(slide).fill_reading(title, body)
    


############################################################################################################
###
### Slide creator
###
############################################################################################################

template_file = 'template.pptx'

# Open the template file
# Create new class instance (prs)
prs = Presentation(template_file)


## Currently trying to sort out adding catechism pages
for i in range (1,23):
    try:
        slide_writer(i, prs)
        print(f"Slide {i} Complete")
    except:
        print(f"Error on slide {i}")

filename = f"{cycle[0]}_week_{datetime.now().isocalendar()[1]}.pptx"
file_path = input("Enter a file path: ")

if os.path.exists(file_path):
    print(f"{file_path} is a valid file path")
    file_path = os.path.join(file_path, filename)
else:
    print(f"{file_path} is not a valid file path")
    try:
        desktop = os.path.join(os.environ["USERPROFILE"], "Desktop")
    except:
        desktop = os.path.join(os.environ["HOME"], "OneDrive\\Desktop")
    file_path = os.path.join(desktop, filename)

# Save the PowerPoint file
prs.save(file_path)

input("Press enter to exit")

### To create the exe version of this file run this command in the terminal
### pyinstaller --onefile --add-data "template.pptx;." --add-data "components.json;." --add-data "psalms.json;." --add-data "wsc.json;." full_code.py
