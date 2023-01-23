import re
import re
import json
import html
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

#############################################################
## Creates the American to British ESV Dictionary
list1 = [
         "(\W)([Ss])avior(s)?(\W)",
         "(\W)neighbor(s|ing)?(\W)",
         "(\W)favor(able|ite|s|ed|itism)?(\W)",
         "(\W)Favor(\W)",
         "(\W)labor(ed|s)?(\W)",
         "(\W)(vap|vig)or(\W)",
         "(\W)clamor(\W)",
         "(\W)([Ss])plendor(\W)",
         "(\W)color(s|ed)?(\W)",
         "(\W)([Hh])onor(s|able|ing|ed)?(\W)",
         "(\W)dishonor(s|able|ing|ed)?(\W)",
         "(\W)travel(ed|er|ers|ing)(\W)",
         "(\W)marvel(ous|ously|ed|ing)(\W)",
         "(\W)([Cc])ounsel(or|ors|ed)(\W)",
         "(\W)plow(s|ed|ers|ing|man|men|share|shares)?(\W)",
         "(\W)judgment(s)?(\W)",
         "(\W)(recogn|[Rr]eal|[Oo]rgan|[Ss]ymbol|bapt|critic|apolog|sympath)iz(e|ed|es|ing)?(\W)",
         "(\W)(un)?authorized(\W)",
         "(\W)(centi)?meters(\W)",
         "(\W)liter(s)?(\W)",
         "(\W)scepter(s)?(\W)",
         "(\W)worship(ed|er|ers|ing)(\W)",
         "(\W)quarrel(ed|ing)(\W)",
         "(\W)benefited(\W)",
         "(\W)signaled(\W)",
         "(\W)paralyzed(\W)",
         "(\W)fulfill(s|ment)?(\W)",
         "(\W)skillful(ly)?(\W)",
         "(\W)jewelry(\W)",
         "(\W)(De|de|of)fense(s|less)?(\W)",
         "(\W)([Ss])ulfur?(\W)"
    ]

list2 = [
        r'\1\2aviour\3\4',
        r'\1neighbour\2\3',
        r'\1favour\2\3',
        r'\1Favour\2',
        r'\1labour\2\3',
        r'\1\2our\3',
        r'\1clamour\2',
        r'\1\2plendour\3',
        r'\1colour\2\3',
        r'\1\2onour\3\4',
        r'\1dishonour\2\3',
        r'\1travell\2\3',
        r'\1marvell\2\3',
        r'\1\2ounsell\3\4',
        r'\1plough\2\3',
        r'\1judgement\2\3',
        r'\1\2is\3\4',
        r'\1\2authorised\3',
        r'\1\2metres\3',
        r'\1litre\2\3',
        r'\1sceptre\2\3',
        r'\1worshipp\2\3',
        r'\1quarrell\2\3',
        r'\1benefitted\2',
        r'\1signalled\2',
        r'\1paralysed\2',
        r'\1fulfil\2\3',
        r'\1skilful\2\3',
        r'\1jewellery\2',
        r'\1\2fence\3\4',
        r'\1\2ulphur\3'
    ]
#############################################################


##############################################################
## Data Loading Section

## Generates list to use in below function
anglo_list = list(zip(list1, list2))

## Need to convert to function
def create_json_var(name):
    with open(f"{name}.json", "r") as f:
        value = json.load(f)
        globals()[name + "_json"] = value

create_json_var("components")
create_json_var("psalms")
create_json_var("wsc")

API_KEY = 'ade14fe748fbb522b8dfb225ec6b222fa148cddc'
API_URL = 'https://api.esv.org/v3/passage/text/'
##############################################################
  
## Extracts the relevant details from the components JSON
def component_assigner(val, flag = None):
    if flag != 'communion':
        res = list(filter(lambda item: item["Component"] == val, components_json))
    else:
        initial = list(filter(lambda item: item["Component"] == "communion", components_json))
        res = list(filter(lambda item: item["Component"] == val, initial[0]["Content"]))
        
    speaker = res[0]["Content"]["speaker"]
    return_list = []
    for i in range(1,5):
        try:
            return_list.append(res[0]["Content"][str(i)])
        except:
            pass
    return speaker, return_list
'''
Returns the speaker, which most likely goes on every slide
And returns the text for each slide, in a list, with the content for each slide in each element
'''

## Function to extract the appropriate catechism question
def catechism_finder(val):
    real_val = int(val) - 1
    question = wsc_json['Data'][real_val]["Question"]
    answer = wsc_json['Data'][real_val]["Answer"]
    return (question,answer, val)


## Function to convert American language text to British
def replace_text(text, rep_text=anglo_list):
    for pattern, replacement in rep_text:
        pattern = re.compile(pattern)
        text = pattern.sub(replacement, text)
    return text

## Function to convert digits in verse to superscript
def super_verse(text):
    pattern = re.compile(r'\[([^\[\]]*)\]\s')
    matches = pattern.findall(text)
    for match in matches:
        superscript = ''.join(['⁰¹²³⁴⁵⁶⁷⁸⁹'[int(x)] if x.isdigit() else x for x in match])
        text = text.replace('[' + match + '] ', superscript)
    return text
    

## ESV API to return text from passage text
def get_esv_text(passage):
  params = {
    'q': passage,
    'include-headings': False,
    'include-footnotes': False,
    'include-verse-numbers': True,
    'include-short-copyright': False,
    'include-passage-references': False
  }
  
  headers = {
    'Authorization': 'Token %s' % API_KEY
  }
  
  response = requests.get(API_URL, params=params, headers=headers)
  
  passages = response.json()['passages']

  result = passages[0].strip()
  result = super_verse(result)
  result = replace_text(result)
  
  if passages:
    return (result,passage)
  else:
    return 'Error: Passage not found'

'''  
result = get_esv_text("John 3:16-19")
print(result)
'''

## Hymn Scraper Class
class HymnScraper:
    def __init__(self, url1 = None, url2 = None):
        self.url1 = url1
        self.url2 = url2
        self.title = None
        self.author = None
        self.copyr = None
        self.n_verses = None
        self.composer = None
        self.tune = None
        self.meter = None
        
    def hymn_lyrics(self):
        ## Word Extraction
        response = requests.get(self.url1)
        soup = BeautifulSoup(response.content, 'html.parser')
        try:
            self.title = soup.find("span", class_="hy_infoLabel", string="Title:").find_next().text
        except:
            None
        try:
            self.author = soup.find("span", class_="hy_infoLabel", string="Author:").find_next().text
        except:
            try:
                self.author = soup.find("span", class_="hy_infoLabel", string="Author (attributed to):").find_next().text
                self.author = self.author + " (atrb)"
            except:
                self.author = None
        try:
            self.copyr = soup.find("span", class_="hy_infoLabel", string="Copyright:").find_next().text
        except:
            self.copyr = None

        try:
            body_text = soup.select_one('div#at_fulltext.authority_section div div.authority_columns')
            paragraphs = body_text.find_all('p')
            verses = []
            for p in paragraphs:
                verses.append(p.text)

            self.n_verses = []
            for item in verses:
                if item[0].isdigit():
                    items = item.split(" ")
                    item = " ".join(items[1:])
                    if item.rstrip().endswith(" [Refrain]"):
                        item = item.split(" [Refrain]")[0]
                        self.n_verses.append(item)
                        self.n_verses.append(refrain)
                    else:
                        self.n_verses.append(item)
                if item.split(":")[0] == 'Refrain':
                    refrain = item.split(":")[1] + "\n"
                    self.n_verses.append(refrain)
        except:
            self.n_verses = None

    def tune_details(self):
        ## Tune Extraction
        response = requests.get(self.url2)
        soup = BeautifulSoup(response.content, 'html.parser')
        try:
            self.composer = soup.find("span", class_="hy_infoLabel", string="Composer:").find_next().text
        except:
            try:
                self.composer = soup.find("span", class_="hy_infoLabel", string="Composer (attributed to):").find_next().text
                self.composer = self.composer + " (atrb)"
            except:
                self.composer = "Unknown"
        try:
            self.tune = soup.find("span", class_="hy_infoLabel", string="Title:").find_next().text
        except:
            self.tune = None
        try:
            self.meter = soup.find("span", class_="hy_infoLabel", string="Meter:").find_next().text
        except:
            self.meter = None


    def get_lyrics(self):
        return (self.n_verses,self.title, self.author, self.composer, self.tune, self.copyr)
    
    def get_tune(self):
        return (self.composer, self.tune, self.meter)

## Hymn Scraper Functions
## Function to instantiate class and extract details
def song_details(url1,url2):
    scraper = HymnScraper(url1, url2)
    scraper.hymn_lyrics()
    scraper.tune_details()
    return scraper.get_lyrics()

def tune_details(url2):
    scraper = HymnScraper(None, url2)
    scraper.tune_details()
    return scraper.get_tune()

# function to convert to superscript
def get_super(x):
    normal = "0123456789"
    super_s = "⁰¹²³⁴⁵⁶⁷⁸⁹"
    res = x.maketrans(''.join(normal), ''.join(super_s))
    return x.translate(res)

def psalm_getter(psalm_name, psalms_json=psalms_json):
    
    match = re.search(r"^(Psalm )?(\d{1,3})(:(([1-9]\d{0,2})-([1-9]\d{0,2})))?(\s\(([a-zA-Z])\))?(\s\((\d{0,2})\))?$", psalm_name)

    if match:
        if match.group(2):
            psalm = match.group(2)
        if match.group(7):
            version = match.group(8)
        else:
            version = "a"
        if match.group(9):
            section = match.group(10)
        else:
            section = None
        if match.group(4):
            verses = match.group(4)
        elif section != None:
            verses = "1-300"
        else:
            verses = "1-30"

    if section == None:
        body = [d for d in psalms_json if d["Psalm"] == str(psalm) and d["Content"]["Version"] == version][0]['Content']['Body']
        meter = [d for d in psalms_json if d["Psalm"] == str(psalm) and d["Content"]["Version"] == version][0]['Content']['Meter']
    else:
        body = [d for d in psalms_json if d["Psalm"] == str(psalm) and d["Content"]["Section"] == str(section)][0]['Content']['Body']
        meter = [d for d in psalms_json if d["Psalm"] == str(psalm) and d["Content"]["Section"] == str(section)][0]['Content']['Meter']

    lst = []
    rng = [int(n) for n in verses.split("-")]
    for ints in range(rng[0],rng[1]):
        try:
            lst.append(body[str(ints)])
        except:
            pass

    lst = list(map(get_super,lst))

    return (psalm_name, meter,lst)

'''
# Pass in the string of the Psalm in the below format
##string = "Psalm 19:2-10 (a)"
# Swapping a either for section (number) or version (letter)
psalm, meter, lst = psalm_getter(string, psalms_json)
print(psalm,meter,lst)
'''

class SlideFill:
    def __init__(self, slide):
        self.slide = slide
    
    def fill_main(self, title, body, copy):
        for idx, placeholder in enumerate(self.slide.placeholders):
            try:
                placeholder.text_frame.clear()
                if idx == 0:
                    new_body = title
                elif idx == 1:
                    new_body = body
                elif idx == 2:
                    new_body = copy
                placeholder.text_frame.text = new_body
            except:
                pass
    
    def fill_component(self, title, body, address):
        for idx, placeholder in enumerate(self.slide.placeholders):
            try:
                placeholder.text_frame.clear()
                if idx == 0:
                    placeholder.text_frame.text = title
                elif idx == 1:
                    p = placeholder.text_frame.add_paragraph()
                    run = p.add_run()
                    run.font.size = Pt(28)
                    run.font.color.rgb = RGBColor(161, 38, 38)
                    run.text = f"{address} "
                    run = p.add_run()
                    run.font.size = Pt(28)
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    run.text = body    
            except:
                pass
    
    def fill_psalm(self, title, body, copy, meter):
        for idx, placeholder in enumerate(self.slide.placeholders):
            try:
                placeholder.text_frame.clear()
                if idx == 0:
                    new_body = title
                elif idx == 1:
                    new_body = body
                elif idx == 2:
                    new_body = copy
                elif idx == 3:
                    new_body = meter
                placeholder.text_frame.text = new_body
            except:
                pass

def flag_format(flag):
    try:
        ## Such as 'notices' and 'song1'
        if type(flag) == str:
            return globals()[flag]
        elif type(flag) == list:
            try:
                ## For 'component versions
                if type(flag) == list:
                    return component_assigner(flag[0])
            except:
                ## For 'call to worship'
                return (globals()[flag[0]],flag[1])
    except:
        return "Error"

def slide_maker(layout_type):
    slide_layout = prs.slide_layouts[layout_type]
    slide = prs.slides.add_slide(slide_layout)
    return slide


def slide_writer(flag):
    # Populate the placeholders on the slide with data from variables
    flag_val = flag_format(slide_dict[flag])

    
    if flag == 3:
        slide = slide_maker(1)
        # This works for call to worship
        title = flag_val[1]
        body = flag_val[0][0]
        copy = flag_val[0][1]
        SlideFill(slide).fill_main(title, body, copy)
    elif flag in song_list:
        # This works for psalms
        try:
            flag_val[0].find("Psalm") != -1
            list_len = len(flag_val[2])
            for i in range(0,list_len):
                slide = slide_maker(3)
                title = flag_val[0]
                body = flag_val[2][i]
                copy = f"Words: Sing Psalms! © 2003 Free Church of Scotland\nComposer: {psalm_tune[0]}\nTune: {psalm_tune[1]}\n©: Public Domain\nCCLI: 522221"
                meter = f"Meter: {flag_val[1]}"
                SlideFill(slide).fill_psalm(title, body, copy, meter)
        except:
            # This works for songs
            try:
                list_len = len(flag_val[0])
                for i in range(0,list_len):
                    slide = slide_maker(0)
                    title = flag_val[1]
                    body = flag_val[0][i]
                    # Could put this outside the loop, in the future to only fill on the final slide
                    copy = f"Words: {flag_val[2]}\nComposer: {flag_val[3]}\nTune: {flag_val[4]}\n©: {flag_val[5]}\nCCLI: 522221"
                    SlideFill(slide).fill_main(title, body, copy)
            except:
                slide = slide_maker(0)
                try:
                    title = flag_val[1]
                except:
                    title = slide_dict[flag]
                body = "Not in Public Domain"
                try:
                    copy = f"Words: {flag_val[2]}\nComposer: {flag_val[3]}\nTune: {flag_val[4]}\n©: {flag_val[5]}\nCCLI: 522221"
                except:
                    copy = "Error"
                SlideFill(slide).fill_main(title, body, copy)
                
    elif flag in component_list:
        # This works for components
        list_len = len(flag_val[1])
        for i in range(0,list_len):
            slide = slide_maker(2)
            title = slide_dict[flag][1]
            address = flag_val[0]
            body = flag_val[1][i]
            SlideFill(slide).fill_component(title, body, address)
    