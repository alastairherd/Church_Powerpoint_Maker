import importlib.util
import sys
import pandas as pd
import requests
import variables
from pptx import Presentation


def load_functions(monkeypatch, resp=None):
    """Import functions module with patched dependencies."""
    if resp is None:
        class DummyResp:
            def json(self):
                return {'passages': ['']}
        resp = DummyResp()

    monkeypatch.setattr(pd, 'read_csv', lambda *a, **k: pd.DataFrame({'Service Column':[None]*20}))
    monkeypatch.setattr(requests, 'get', lambda *a, **k: resp)

    spec = importlib.util.spec_from_file_location('functions', 'functions.py')
    module = importlib.util.module_from_spec(spec)
    module.var = variables
    module.anglo_list = variables.anglo_list
    spec.loader.exec_module(module)
    return module


def test_replace_text(monkeypatch):
    funcs = load_functions(monkeypatch)
    text = ' my neighbor loves color '
    result = funcs.replace_text(text, funcs.var.anglo_list)
    assert 'neighbour' in result
    assert 'colour' in result


def test_super_verse(monkeypatch):
    funcs = load_functions(monkeypatch)
    result = funcs.super_verse('text [123] more')
    assert '¹²³' in result


def test_psalm_getter(monkeypatch):
    funcs = load_functions(monkeypatch)
    psalm, meter, verses = funcs.psalm_getter('Psalm 119 (12)')
    assert meter == '8 8 8 8 8 8'
    assert verses[0].startswith('⁸⁹Eternal is your word')


def test_catechism_finder(monkeypatch):
    funcs = load_functions(monkeypatch)
    q, a, num = funcs.catechism_finder(1)
    assert 'chief end of man' in q
    assert 'glorify God' in a


def test_component_assigner(monkeypatch):
    funcs = load_functions(monkeypatch)
    speaker, items = funcs.component_assigner('confession')
    assert speaker == 'All.'
    assert isinstance(items, list) and len(items) > 0


def test_get_esv_text(monkeypatch):
    class DummyResp:
        def json(self):
            return {'passages': [' my neighbor [12] test']}
    funcs = load_functions(monkeypatch, DummyResp())
    result, ref = funcs.get_esv_text('John 3:16')
    assert 'neighbour' in result
    assert '¹²' in result
    assert ref == 'John 3:16'


def test_slidefill(monkeypatch):
    funcs = load_functions(monkeypatch)
    sys.modules['functions'] = funcs
    import slide_making
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide_making.SlideFill(slide).fill_main('Title', 'Body', 'Copy')
    texts = [ph.text for ph in slide.placeholders][:3]
    assert texts == ['Title', 'Body', 'Copy']
